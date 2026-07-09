from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"C:\Users\HARSH TIWARI\Documents\Codex\2026-07-02\h\outputs\chatbot-ui")
ASSETS = ROOT / "frontend" / "assets"
OUT = ROOT / "Byte-Bot_Project_CPU_Format.pptx"
TMP = ROOT / "miscellaneous" / "ppt_assets"
TMP.mkdir(parents=True, exist_ok=True)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(9, 34, 74)
DARK = RGBColor(6, 18, 45)
BLUE = RGBColor(30, 84, 148)
LIGHT = RGBColor(242, 247, 255)
TEXT = RGBColor(30, 42, 60)
ACCENT = RGBColor(255, 198, 28)
CYAN = RGBColor(38, 166, 220)
GREEN = RGBColor(50, 170, 110)
RED = RGBColor(210, 65, 65)


def convert_image(src_name, out_name=None):
    src = ASSETS / src_name
    if not src.exists():
        return None
    out = TMP / (out_name or f"{src.stem}.png")
    try:
        img = Image.open(src)
        if getattr(img, "is_animated", False):
            img.seek(0)
        img.convert("RGBA").save(out)
        return out
    except Exception:
        return src


def add_bg(slide, color=RGBColor(255, 255, 255)):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()


def add_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.68))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(0.12), Inches(10.5), Inches(0.45))
    p = tx.text_frame.paragraphs[0]
    p.text = title.upper()
    p.font.name = "Cambria"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    if subtitle:
        st = slide.shapes.add_textbox(Inches(9.4), Inches(0.18), Inches(3.45), Inches(0.35))
        p2 = st.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(210, 225, 245)
        p2.alignment = PP_ALIGN.RIGHT


def add_footer(slide):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.1), SLIDE_W, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(6.5), Inches(7.17), Inches(6.35), Inches(0.22))
    p = tx.text_frame.paragraphs[0]
    p.text = "Byte-Bot AI Chatbot Project | Frontend + FastAPI + LangChain + MongoDB"
    p.font.name = "Calibri"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(235, 245, 255)
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title, subtitle=None):
    t = slide.shapes.add_textbox(Inches(0.55), Inches(0.95), Inches(6.8), Inches(0.55))
    p = t.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Cambria"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = NAVY
    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.58), Inches(1.48), Inches(6.9), Inches(0.36))
        ps = s.text_frame.paragraphs[0]
        ps.text = subtitle
        ps.font.name = "Calibri"
        ps.font.size = Pt(14)
        ps.font.color.rgb = BLUE


def add_bullets(slide, x, y, w, h, bullets, font_size=16, title=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    first = True
    if title:
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Cambria"
        p.font.bold = True
        p.font.size = Pt(font_size + 2)
        p.font.color.rgb = NAVY
        first = False
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(6)


def add_section_card(slide, x, y, w, h, heading, lines, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT
    card.line.color.rgb = RGBColor(185, 205, 230)
    card.line.width = Pt(1.4)
    hbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    hbox.fill.solid()
    hbox.fill.fore_color.rgb = accent
    hbox.line.fill.background()
    th = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.05), Inches(w - 0.25), Inches(0.22))
    p = th.text_frame.paragraphs[0]
    p.text = heading
    p.font.name = "Cambria"
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(255, 255, 255)
    add_bullets(slide, x + 0.18, y + 0.48, w - 0.35, h - 0.58, lines, font_size=12)


def add_image_fit(slide, path, x, y, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_flow(slide, labels, x=0.75, y=2.0, box_w=1.55, box_h=0.72, gap=0.32):
    palette = [BLUE, CYAN, GREEN, ACCENT, RED, NAVY]
    for i, label in enumerate(labels):
        bx = x + i * (box_w + gap)
        sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(bx), Inches(y), Inches(box_w), Inches(box_h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = palette[i % len(palette)]
        sh.line.color.rgb = RGBColor(20, 30, 45)
        p = sh.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(20, 20, 20) if i == 3 else RGBColor(255, 255, 255)
        if i:
            line = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(bx - gap + 0.03),
                Inches(y + box_h / 2),
                Inches(bx - 0.03),
                Inches(y + box_h / 2),
            )
            line.line.color.rgb = NAVY
            line.line.width = Pt(2)


def build_deck():
    mascot = convert_image("landing-mascot.webp")
    chat_banner = ASSETS / "chat-sunset-banner.png"
    web_icon = ASSETS / "web-search-tool.png"
    upload_icon = ASSETS / "upload-toolkit.png"
    voice_icon = ASSETS / "voice-output-tool.png"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_image_fit(s, chat_banner, 0, 0, 13.333, 3.05)
    add_header(s, "BYTE-BOT AI CHATBOT", "Project Presentation")
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.25), Inches(6.3), Inches(2.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.fill.transparency = 8
    box.line.color.rgb = RGBColor(210, 220, 235)
    t = s.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(5.8), Inches(0.7))
    p = t.text_frame.paragraphs[0]
    p.text = "Byte-Bot"
    p.font.name = "Cambria"
    p.font.bold = True
    p.font.size = Pt(44)
    p.font.color.rgb = NAVY
    st = s.shapes.add_textbox(Inches(0.92), Inches(2.25), Inches(5.8), Inches(0.55))
    p = st.text_frame.paragraphs[0]
    p.text = "A Pixel-Style AI Chatbot With Backend Memory, Tools, Voice Roadmap And Deployment Plan"
    p.font.name = "Calibri"
    p.font.size = Pt(15)
    p.font.color.rgb = TEXT
    info = s.shapes.add_textbox(Inches(0.92), Inches(3.02), Inches(5.9), Inches(0.8))
    tf = info.text_frame
    for idx, line in enumerate([
        "Submitted By: Harsh Tiwari",
        "Guided By: Faculty / Department",
        "Stack: HTML, CSS, JavaScript, FastAPI, LangChain, Groq, MongoDB",
    ]):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Calibri"
        p.font.size = Pt(12)
        p.font.color.rgb = TEXT
    add_image_fit(s, mascot, 8.55, 3.65, 2.15, 2.15)
    add_footer(s)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "TABLE OF CONTENT")
    add_footer(s)
    items = [
        "Introduction", "Current Gap and Vision", "Objectives", "Problem Statement",
        "Solution / Proposed System", "System Architecture", "Technologies Used",
        "Features / Modules", "Frontend and Backend Workflow", "Integration and Deployment",
        "Future Enhancements", "Conclusion and Thank You",
    ]
    for i, item in enumerate(items):
        col = 0 if i < 6 else 1
        row = i if i < 6 else i - 6
        x, y = 1.0 + col * 5.9, 1.25 + row * 0.78
        num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.36), Inches(0.36))
        num.fill.solid()
        num.fill.fore_color.rgb = ACCENT
        num.line.color.rgb = NAVY
        p = num.text_frame.paragraphs[0]
        p.text = str(i + 1)
        p.font.bold = True
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
        tx = s.shapes.add_textbox(Inches(x + 0.48), Inches(y - 0.01), Inches(4.9), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = item
        p.font.name = "Calibri"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT

    slide_data = [
        ("INTRODUCTION", "Background", "Why Byte-Bot was built", [
            "Many chatbot projects stop at a simple textbox and do not feel interactive or memorable.",
            "Byte-Bot began as a pixel-style frontend UI, then grew into a real AI chatbot connected to a backend.",
            "The goal is to make the bot useful and charming: it can chat, remember sessions, search the web, and later understand user documents.",
            "The project focuses on beginner-readable code, clear documentation, and deployment readiness.",
        ]),
    ]

    for header, title, subtitle, bullets in slide_data:
        s = prs.slides.add_slide(blank)
        add_bg(s)
        add_header(s, header)
        add_footer(s)
        add_title(s, title, subtitle)
        add_bullets(s, 0.85, 1.75, 6.35, 4.8, bullets, 16)
        add_image_fit(s, mascot, 8.4, 2.1, 2.2, 2.2)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "INTRODUCTION")
    add_footer(s)
    add_title(s, "Current Gap and Vision")
    add_section_card(s, 0.8, 1.75, 5.55, 4.45, "Current Gap", [
        "Frontend-only chat demos cannot answer intelligently.",
        "Plain LLM wrappers forget previous messages without memory.",
        "Users need search, file understanding, and voice features in one simple interface.",
        "Deployment becomes hard when frontend, backend, database, and API keys are not organized.",
    ], RED)
    add_section_card(s, 6.85, 1.75, 5.55, 4.45, "Vision", [
        "Build a scalable AI companion with a polished pixel UI.",
        "Connect the UI to FastAPI, LangChain, Groq, MongoDB, and tool modules.",
        "Support future RAG over PDFs, CSVs, text files, and images.",
        "Keep the project understandable enough to explain in viva or classroom review.",
    ], GREEN)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "OBJECTIVES")
    add_footer(s)
    add_title(s, "Main Objectives")
    add_section_card(s, 0.75, 1.65, 3.85, 4.75, "User Focus", [
        "Send messages from a friendly pixel UI.",
        "Receive short, helpful, personality-rich answers.",
        "Use buttons for upload, voice, search, and reset.",
    ], BLUE)
    add_section_card(s, 4.75, 1.65, 3.85, 4.75, "System Focus", [
        "Use FastAPI for clean backend endpoints.",
        "Store users, sessions, and memory in MongoDB.",
        "Keep API keys safely in environment variables.",
    ], CYAN)
    add_section_card(s, 8.75, 1.65, 3.85, 4.75, "Functional Focus", [
        "Integrate Groq LLM through LangChain.",
        "Add document search using chunking and retrieval.",
        "Prepare for TTS, STT, hosting, and future tools.",
    ], GREEN)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "PROBLEM STATEMENT")
    add_footer(s)
    add_title(s, "Problem Statement")
    add_bullets(s, 0.85, 1.75, 6.35, 4.8, [
        "A basic chatbot UI is not enough because it cannot preserve memory or call real backend logic.",
        "Without backend organization, login, database storage, and deployment become confusing.",
        "Search results can be messy unless the bot summarizes context first and keeps links separate.",
        "Uploaded files need chunking and retrieval before the bot can answer from PDF/CSV/TXT content.",
        "Voice and assistant tools need a clear roadmap so they do not break the stable chat system.",
    ], 16)
    add_section_card(s, 7.45, 2.0, 4.7, 1.8, "Core Challenge", [
        "Merge many parts into one stable beginner-friendly product without losing the Byte-Bot personality."
    ], ACCENT)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "SOLUTION / PROPOSED SYSTEM")
    add_footer(s)
    add_title(s, "Proposed System")
    solutions = [
        ("Pixel Chat Interface", "HTML, CSS, and JavaScript create the visual chat experience."),
        ("FastAPI Backend", "Python routes receive messages, reset memory, and expose tools."),
        ("AI Brain", "LangChain sends the prompt and history to Groq."),
        ("MongoDB Memory", "Registration, login, and chat sessions can be stored in backend database."),
        ("Tool Layer", "Search, upload, voice, and future assistant tools are added as separate modules."),
    ]
    for i, (heading, line) in enumerate(solutions):
        add_section_card(s, 0.75, 1.45 + i * 0.86, 6.05, 0.72, heading, [line], [BLUE, CYAN, GREEN, ACCENT, RED][i])
    add_image_fit(s, upload_icon, 7.5, 1.6, 1.45, 1.45)
    add_image_fit(s, web_icon, 9.05, 1.6, 1.45, 1.45)
    add_image_fit(s, voice_icon, 10.6, 1.6, 1.45, 1.45)
    add_bullets(s, 7.55, 3.5, 4.6, 2.1, [
        "The base chatbot stays stable first; advanced abilities are added step by step.",
        "Each tool is isolated so one broken feature does not break the whole chat flow.",
    ], 15)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "SYSTEM ARCHITECTURE")
    add_footer(s)
    add_title(s, "Architecture Flow")
    add_flow(s, ["User", "Frontend", "FastAPI", "LangChain", "Groq", "MongoDB", "Reply"], x=0.65, y=2.0, box_w=1.4, box_h=0.82, gap=0.26)
    add_bullets(s, 0.85, 3.45, 11.4, 2.1, [
        "Frontend sends JSON data such as message and session_id using fetch().",
        "FastAPI receives the request, loads memory, adds prompt instructions, and calls the AI model.",
        "MongoDB stores user/login data and can store chat/session data for scalable memory.",
        "The backend returns JSON, and JavaScript renders it inside the pixel chat bubble.",
    ], 16)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "TECHNOLOGIES USED")
    add_footer(s)
    add_title(s, "Tools Used and Why")
    tech = [
        ("HTML / CSS / JS", "Simple frontend stack, easy to edit, good for learning web basics."),
        ("FastAPI", "Modern Python framework, fast JSON APIs, automatic docs, cleaner than Flask for API work."),
        ("LangChain + Groq", "LangChain organizes prompts and memory; Groq gives fast LLM responses."),
        ("MongoDB Atlas", "Cloud database for registration, login, sessions, and future scalable storage."),
        ("SerpApi", "Search API for web/image results when user asks for current external information."),
        ("Vercel / Render", "Free-friendly hosting path: frontend on Vercel, backend on Render/Railway."),
    ]
    for i, (heading, line) in enumerate(tech):
        add_section_card(s, 0.75 + (i % 2) * 6.1, 1.45 + (i // 2) * 1.55, 5.65, 1.15, heading, [line], [BLUE, CYAN, GREEN, ACCENT, RED, NAVY][i])

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "FEATURES / MODULES")
    add_footer(s)
    add_title(s, "Current and Planned Modules")
    add_section_card(s, 0.75, 1.45, 3.8, 4.85, "Completed / Started", [
        "Pixel login and chat UI", "Backend /chat endpoint", "Groq + LangChain reply",
        "MongoDB login connection", "Web search and image search", "Upload chunking prototype",
    ], GREEN)
    add_section_card(s, 4.75, 1.45, 3.8, 4.85, "In Progress", [
        "Cleaner RAG answers", "PDF/CSV/TXT semantic retrieval", "Better response formatting",
        "Deployment configuration", "Teacher-ready documentation",
    ], CYAN)
    add_section_card(s, 8.75, 1.45, 3.8, 4.85, "Future", [
        "Voice input using Whisper", "Voice output using TTS", "Assistant tools/actions",
        "Long-term memory", "Authentication polish", "Production monitoring",
    ], ACCENT)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "WORKFLOW AND INTEGRATION")
    add_footer(s)
    add_title(s, "How Everything Was Merged")
    add_bullets(s, 0.85, 1.55, 5.9, 4.9, [
        "Designed the pixel frontend with static HTML, CSS animations, mascot, and chat bubbles.",
        "Added JavaScript event handling for Send, Reset, Upload, Search, and Voice buttons.",
        "Connected fetch() calls to FastAPI routes such as /chat, /reset, and tool endpoints.",
        "Backend loads environment variables, database connection, prompt, memory, and tool modules.",
        "Responses are returned as JSON and rendered back into the chat UI.",
    ], 15)
    add_flow(s, ["Click", "JS", "Fetch", "API", "AI/DB", "Bubble"], x=7.0, y=2.0, box_w=1.0, box_h=0.72, gap=0.2)
    add_section_card(s, 7.0, 3.35, 5.2, 1.65, "Important Integration Rule", [
        "Frontend knows the backend URL only. Secrets, AI keys, database logic, and tool logic stay in backend."
    ], RED)

    s = prs.slides.add_slide(blank)
    add_bg(s)
    add_header(s, "FUTURE ENHANCEMENTS")
    add_footer(s)
    add_title(s, "Roadmap")
    road = [
        ("TTS", "Browser speechSynthesis first; later ElevenLabs/OpenAI TTS."),
        ("STT", "Whisper/OpenAI transcription endpoint or browser speech recognition."),
        ("RAG", "Upload PDFs, CSVs, TXT files, split into chunks, retrieve relevant context."),
        ("Tools", "Search, calculator, notes, image search, and personal assistant actions."),
        ("Deployment", "Vercel/Netlify frontend plus Render/Railway backend and MongoDB Atlas."),
        ("Security", "Hash passwords, add tokens, protect environment variables, rate limit APIs."),
    ]
    for i, (heading, line) in enumerate(road):
        add_section_card(s, 0.75 + (i % 2) * 6.1, 1.45 + (i // 2) * 1.55, 5.65, 1.12, heading, [line], [BLUE, CYAN, GREEN, ACCENT, RED, NAVY][i])

    s = prs.slides.add_slide(blank)
    add_bg(s, DARK)
    add_image_fit(s, chat_banner, 0, 0, 13.333, 2.55)
    add_header(s, "CONCLUSION AND THANK YOU")
    add_footer(s)
    add_section_card(s, 0.9, 1.55, 6.1, 3.8, "Conclusion", [
        "Byte-Bot has moved from a frontend demo toward a real AI assistant project.",
        "The system combines pixel UI, FastAPI backend, AI model calls, memory, search, and upload-tool planning.",
        "The chosen tools are beginner-friendly, scalable, and suitable for local testing and future hosting.",
        "Next focus: stable deployment, better RAG, voice features, and secure authentication.",
    ], GREEN)
    add_image_fit(s, mascot, 8.1, 2.2, 2.0, 2.0)
    t = s.shapes.add_textbox(Inches(7.45), Inches(4.45), Inches(4.9), Inches(0.8))
    p = t.text_frame.paragraphs[0]
    p.text = "THANK YOU\nAny Questions?"
    p.font.name = "Cambria"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
